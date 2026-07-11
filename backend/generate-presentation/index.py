"""
Генерация PPTX-презентации по теме урока через ИИ-движок. v2
POST / body: {topic, description, slidesCount, audience?, teacherName, teacherSchool}
Возвращает: {pptx_b64, filename, outline}.

Особенности:
- Контент строго по ФГОС и программе Минпросвещения РФ
- Фотографии по теме на каждом слайде (Wikimedia Commons / Unsplash)
- Двухколоночный layout: текст слева, фото справа
- Таймаут 320 секунд, мощная модель ИИ-движка
- Слайд "Содержание" после титульного

GET /?action=ping — проверка доступности ИИ-движка
"""
import json
import os
import io
import re
import time
import base64
import random
import urllib.request
import urllib.error
from datetime import datetime

AUTH_URL = os.environ.get("AUTH_FUNCTION_URL", "https://functions.poehali.dev/b08ae7cf-6c0b-4178-acc9-4b62b2c2a61b")

TOKENS_COST_PRESENTATION = 4000


def spend_ai_tokens(login: str, amount: int, action_label: str = "Презентация") -> tuple[bool, str, float, float]:
    if not login:
        return True, "", 0.0, 0.0
    try:
        req = urllib.request.Request(
            f"{AUTH_URL}?action=spend-tokens",
            data=json.dumps({"login": login, "amount": amount, "action_label": action_label}).encode("utf-8"),
            method="POST",
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=10) as r:
            resp = json.loads(r.read().decode())
        spent_rub = float(resp.get("spent_rub") or 0)
        balance_rub = float(resp.get("balance_rub") or 0)
        return True, "", spent_rub, balance_rub
    except urllib.error.HTTPError as e:
        err_body = {}
        try:
            err_body = json.loads(e.read().decode())
        except Exception:
            pass
        if e.code == 402:
            return False, err_body.get("error", "Недостаточно средств"), 0.0, 0.0
        if e.code == 403:
            return False, err_body.get("error", "Для использования ИИ необходима активная подписка."), 0.0, 0.0
        return True, "", 0.0, 0.0
    except Exception:
        return True, "", 0.0, 0.0


def precheck_ai(login: str, est_tokens: int = 3000):
    """Проверяет ДО обращения к ИИ: подписка и достаточный баланс.
    Возвращает (allowed, http_status, error). Без login — разрешаем."""
    if not login:
        return True, 200, ""
    try:
        req = urllib.request.Request(
            f"{AUTH_URL}?action=precheck-ai",
            data=json.dumps({"login": login, "est_tokens": est_tokens}).encode("utf-8"),
            method="POST",
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=10) as r:
            resp = json.loads(r.read().decode())
        return bool(resp.get("allowed")), 200, ""
    except urllib.error.HTTPError as e:
        err_body = {}
        try:
            err_body = json.loads(e.read().decode())
        except Exception:
            pass
        if e.code == 402:
            return False, 402, err_body.get("error", "Недостаточно средств на балансе ИИ. Пополните баланс.")
        if e.code == 403:
            return False, 403, err_body.get("error", "Для использования ИИ необходима активная подписка.")
        return True, 200, ""
    except Exception:
        return True, 200, ""


from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CORS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
    "Access-Control-Allow-Headers": "Content-Type",
}


def upload_pptx_to_s3(pptx_bytes: bytes, filename: str) -> str:
    """Заливает готовый PPTX в S3 и возвращает CDN-ссылку. Пусто при ошибке."""
    try:
        import boto3
        import uuid
        key_id = os.environ.get("AWS_ACCESS_KEY_ID", "")
        secret = os.environ.get("AWS_SECRET_ACCESS_KEY", "")
        if not key_id or not secret:
            return ""
        s3 = boto3.client(
            "s3",
            endpoint_url="https://bucket.poehali.dev",
            aws_access_key_id=key_id,
            aws_secret_access_key=secret,
        )
        safe = re.sub(r"[^a-zA-Z0-9._-]", "_", filename) or "presentation.pptx"
        key = f"presentations/{uuid.uuid4().hex}_{safe}"
        s3.put_object(
            Bucket="files",
            Key=key,
            Body=pptx_bytes,
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        return f"https://cdn.poehali.dev/projects/{key_id}/bucket/{key}"
    except Exception:
        return ""

# ─── ДИЗАЙН-ТЕМЫ ────────────────────────────────────────────────────────────

THEMES = [
    {   # Глубокий океан — синий с бирюзой
        "name": "ocean",
        "bg":        RGBColor(0xF0, 0xF6, 0xFC),
        "title_bg":  RGBColor(0x09, 0x1E, 0x42),
        "accent":    RGBColor(0x09, 0x1E, 0x42),
        "accent2":   RGBColor(0x00, 0xB4, 0xD8),
        "accent3":   RGBColor(0x48, 0xCA, 0xE4),
        "text":      RGBColor(0x0D, 0x1B, 0x2A),
        "muted":     RGBColor(0x55, 0x70, 0x8B),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0x90, 0xC8, 0xF0),
        "stripe":    RGBColor(0x00, 0xB4, 0xD8),
        "card_bg":   RGBColor(0xE0, 0xF2, 0xFE),
        "label":     "УРОК · ПРЕЗЕНТАЦИЯ",
        "layout":    "left_header",   # шапка слева с крупным номером
    },
    {   # Тёмный лес — изумрудный с золотом
        "name": "forest",
        "bg":        RGBColor(0xF0, 0xF7, 0xF1),
        "title_bg":  RGBColor(0x0F, 0x2D, 0x1F),
        "accent":    RGBColor(0x0F, 0x2D, 0x1F),
        "accent2":   RGBColor(0x2D, 0xC6, 0x5F),
        "accent3":   RGBColor(0xD4, 0xA0, 0x17),
        "text":      RGBColor(0x0A, 0x1F, 0x10),
        "muted":     RGBColor(0x4A, 0x6E, 0x55),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xA8, 0xDF, 0xB8),
        "stripe":    RGBColor(0x2D, 0xC6, 0x5F),
        "card_bg":   RGBColor(0xDC, 0xF5, 0xE7),
        "label":     "УРОК · ПРИРОДА И НАУКА",
        "layout":    "split_diagonal",  # диагональная шапка
    },
    {   # Закат истории — тёмно-красный с золотом
        "name": "sunset",
        "bg":        RGBColor(0xFD, 0xF3, 0xE8),
        "title_bg":  RGBColor(0x3D, 0x0C, 0x02),
        "accent":    RGBColor(0x3D, 0x0C, 0x02),
        "accent2":   RGBColor(0xE8, 0x9A, 0x0C),
        "accent3":   RGBColor(0xC0, 0x39, 0x2B),
        "text":      RGBColor(0x1F, 0x0A, 0x04),
        "muted":     RGBColor(0x7A, 0x50, 0x3A),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xF5, 0xCE, 0x8C),
        "stripe":    RGBColor(0xE8, 0x9A, 0x0C),
        "card_bg":   RGBColor(0xFE, 0xE8, 0xCC),
        "label":     "ИСТОРИЯ И КУЛЬТУРА",
        "layout":    "top_banner",  # широкий баннер сверху
    },
    {   # Квантовый сланец — тёмно-фиолетовый с неоном
        "name": "slate",
        "bg":        RGBColor(0xF4, 0xF2, 0xFF),
        "title_bg":  RGBColor(0x1A, 0x13, 0x36),
        "accent":    RGBColor(0x1A, 0x13, 0x36),
        "accent2":   RGBColor(0x7C, 0x3A, 0xFF),
        "accent3":   RGBColor(0x00, 0xE5, 0xFF),
        "text":      RGBColor(0x12, 0x0C, 0x25),
        "muted":     RGBColor(0x6A, 0x60, 0x90),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xC4, 0xB0, 0xFF),
        "stripe":    RGBColor(0x7C, 0x3A, 0xFF),
        "card_bg":   RGBColor(0xEA, 0xE5, 0xFF),
        "label":     "ТОЧНЫЕ НАУКИ",
        "layout":    "sidebar_dark",  # тёмный сайдбар слева
    },
    {   # Лаборатория — тёмный с кораллом и циановым
        "name": "coral",
        "bg":        RGBColor(0xF8, 0xF5, 0xF2),
        "title_bg":  RGBColor(0x18, 0x22, 0x2F),
        "accent":    RGBColor(0x18, 0x22, 0x2F),
        "accent2":   RGBColor(0xFF, 0x5E, 0x4B),
        "accent3":   RGBColor(0x00, 0xD4, 0xC8),
        "text":      RGBColor(0x10, 0x18, 0x22),
        "muted":     RGBColor(0x6A, 0x72, 0x7E),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xF5, 0xBB, 0xB5),
        "stripe":    RGBColor(0xFF, 0x5E, 0x4B),
        "card_bg":   RGBColor(0xFE, 0xEC, 0xEA),
        "label":     "ХИМИЯ И БИОЛОГИЯ",
        "layout":    "top_banner",
    },
    {   # Арктика — холодный белый с синим льдом
        "name": "arctic",
        "bg":        RGBColor(0xF5, 0xF9, 0xFF),
        "title_bg":  RGBColor(0x0A, 0x2A, 0x4A),
        "accent":    RGBColor(0x0A, 0x2A, 0x4A),
        "accent2":   RGBColor(0x4F, 0xC3, 0xF7),
        "accent3":   RGBColor(0xE0, 0xF7, 0xFA),
        "text":      RGBColor(0x08, 0x20, 0x38),
        "muted":     RGBColor(0x50, 0x70, 0x90),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xB3, 0xE5, 0xFC),
        "stripe":    RGBColor(0x4F, 0xC3, 0xF7),
        "card_bg":   RGBColor(0xE1, 0xF5, 0xFE),
        "label":     "ГЕОГРАФИЯ И ПРИРОДА",
        "layout":    "left_header",
    },
    {   # Рассвет — тёплый кремовый с индиго
        "name": "dawn",
        "bg":        RGBColor(0xFE, 0xF9, 0xF0),
        "title_bg":  RGBColor(0x2C, 0x17, 0x6E),
        "accent":    RGBColor(0x2C, 0x17, 0x6E),
        "accent2":   RGBColor(0xFF, 0x8C, 0x42),
        "accent3":   RGBColor(0xFF, 0xD1, 0x66),
        "text":      RGBColor(0x1A, 0x0E, 0x3A),
        "muted":     RGBColor(0x7A, 0x65, 0x90),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xD4, 0xC0, 0xFF),
        "stripe":    RGBColor(0xFF, 0x8C, 0x42),
        "card_bg":   RGBColor(0xFE, 0xF0, 0xDC),
        "label":     "ЛИТЕРАТУРА И ИСКУССТВО",
        "layout":    "split_diagonal",
    },
    {   # Монохром — элегантный чёрно-белый с акцентом
        "name": "mono",
        "bg":        RGBColor(0xFA, 0xFA, 0xFA),
        "title_bg":  RGBColor(0x10, 0x10, 0x10),
        "accent":    RGBColor(0x10, 0x10, 0x10),
        "accent2":   RGBColor(0xE5, 0x3E, 0x3E),
        "accent3":   RGBColor(0xFF, 0xD7, 0x00),
        "text":      RGBColor(0x18, 0x18, 0x18),
        "muted":     RGBColor(0x70, 0x70, 0x70),
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": RGBColor(0xCC, 0xCC, 0xCC),
        "stripe":    RGBColor(0xE5, 0x3E, 0x3E),
        "card_bg":   RGBColor(0xF0, 0xF0, 0xF0),
        "label":     "УРОК · ПРЕЗЕНТАЦИЯ",
        "layout":    "sidebar_dark",
    },
]

THEME_KEYWORDS = {
    "forest":  ["биол", "экол", "природ", "животн", "растен", "лес", "зоол", "ботан", "органи", "генет", "эволюц"],
    "sunset":  ["истор", "литер", "обществ", "война", "революц", "культур", "философ", "социол", "политол", "граждан", "религ"],
    "slate":   ["физик", "матем", "информат", "програм", "алгебр", "геометр", "алгоритм", "электрон", "квант", "механик", "тригон"],
    "coral":   ["хими", "медиц", "здоровь", "биохим", "анатом", "физиол", "реакц", "молекул", "клетк", "микроб"],
    "arctic":  ["геograph", "геогр", "климат", "страно", "материк", "океан", "река", "горн", "атмосфер", "почв", "ландшафт"],
    "dawn":    ["искусств", "живопис", "музык", "архитект", "театр", "кино", "скульптур", "поэз", "роман", "пьес"],
    "mono":    ["правов", "эконом", "финанс", "право", "конституц", "государств", "юрид", "налог"],
}


def pick_theme(topic: str) -> dict:
    t = topic.lower()
    for theme_name, keywords in THEME_KEYWORDS.items():
        for kw in keywords:
            if kw in t:
                return next(th for th in THEMES if th["name"] == theme_name)
    return random.choice(THEMES)


# ─── ИНДИВИДУАЛЬНЫЙ ДИЗАЙН ──────────────────────────────────────────────────
import colorsys

_CUSTOM_LAYOUTS = ["left_header", "split_diagonal", "top_banner", "sidebar_dark"]


def _hsl(h: float, s: float, l: float) -> "RGBColor":
    """h,s,l в диапазоне 0..1 → RGBColor."""
    r, g, b = colorsys.hls_to_rgb(h % 1.0, max(0.0, min(1.0, l)), max(0.0, min(1.0, s)))
    return RGBColor(int(r * 255), int(g * 255), int(b * 255))


def make_custom_theme(seed_text: str = "", variant: int = 0) -> dict:
    """
    Генерирует уникальную современную дизайн-тему НА ЛЕТУ: палитра, вёрстка и
    графическая композиция собираются процедурно во время запроса (не из готовых
    шаблонов). variant позволяет получить другой вариант для той же темы.
    """
    rnd = random.Random((seed_text + "|" + str(variant) + "|" + str(random.random())).encode("utf-8"))

    # Базовый тон «под тему»: завязан на текст темы, но с лёгкой вариацией —
    # биология тяготеет к зелёному, история к тёплому и т.д., но не одинаково.
    topic_hash = sum(ord(c) for c in seed_text) if seed_text else rnd.randint(0, 360)
    base_h = ((topic_hash % 360) / 360.0 + rnd.uniform(-0.06, 0.06)) % 1.0
    acc2_h = (base_h + rnd.choice([0.5, 0.45, 0.55, 0.33, 0.66])) % 1.0  # доп. акцент
    acc3_h = (base_h + rnd.choice([0.08, -0.08, 0.12])) % 1.0

    title_bg = _hsl(base_h, rnd.uniform(0.55, 0.8), rnd.uniform(0.10, 0.16))   # глубокий тёмный
    accent2  = _hsl(acc2_h, rnd.uniform(0.7, 0.92), rnd.uniform(0.50, 0.60))   # яркий неон-акцент
    accent3  = _hsl(acc3_h, rnd.uniform(0.6, 0.85), rnd.uniform(0.55, 0.68))
    bg       = _hsl(base_h, rnd.uniform(0.18, 0.35), rnd.uniform(0.96, 0.985)) # очень светлый фон
    card_bg  = _hsl(base_h, rnd.uniform(0.25, 0.45), rnd.uniform(0.90, 0.95))
    title_sub= _hsl(base_h, rnd.uniform(0.35, 0.6), rnd.uniform(0.78, 0.86))
    text     = _hsl(base_h, rnd.uniform(0.3, 0.5), rnd.uniform(0.08, 0.13))
    muted    = _hsl(base_h, rnd.uniform(0.2, 0.35), rnd.uniform(0.40, 0.50))

    return {
        "name":      "custom",
        "bg":        bg,
        "title_bg":  title_bg,
        "accent":    title_bg,
        "accent2":   accent2,
        "accent3":   accent3,
        "text":      text,
        "muted":     muted,
        "white":     RGBColor(0xFF, 0xFF, 0xFF),
        "title_sub": title_sub,
        "stripe":    accent2,
        "card_bg":   card_bg,
        "label":     "ИНДИВИДУАЛЬНЫЙ ДИЗАЙН",
        "layout":    rnd.choice(_CUSTOM_LAYOUTS),
        "decor":     generate_decor_recipe(rnd),
    }


_THEME_COLOR_KEYS = ["bg", "title_bg", "accent", "accent2", "accent3",
                     "text", "muted", "white", "title_sub", "stripe", "card_bg"]


def theme_to_payload(theme: dict) -> dict:
    """Сериализует тему в JSON (цвета → hex-строки) для передачи outline→build."""
    out = {"name": theme["name"], "label": theme["label"], "layout": theme["layout"],
           "decor": theme.get("decor") or {}}
    for k in _THEME_COLOR_KEYS:
        # RGBColor — подкласс str, str(color) даёт 6-символьный hex (напр. '1A2B3C')
        out[k] = str(theme[k])
    return out


def theme_from_payload(payload: dict) -> dict:
    """Восстанавливает тему из JSON-payload (hex-строки → RGBColor)."""
    theme = {"name": payload.get("name", "custom"),
             "label": payload.get("label", "ИНДИВИДУАЛЬНЫЙ ДИЗАЙН"),
             "layout": payload.get("layout", "top_banner"),
             "decor": payload.get("decor") or {}}
    for k in _THEME_COLOR_KEYS:
        hexv = str(payload.get(k, "FFFFFF")).lstrip("#") or "FFFFFF"
        theme[k] = RGBColor.from_string(hexv)
    return theme


def _resp(status: int, body: dict) -> dict:
    return {
        "statusCode": status,
        "headers": {**CORS, "Content-Type": "application/json"},
        "body": json.dumps(body, ensure_ascii=False),
        "isBase64Encoded": False,
    }


# ─── YANDEXGPT API ────────────────────────────────────────────────────────────

YANDEX_GPT_URL = "https://llm.api.cloud.yandex.net/foundationModels/v1/completion"


def _yandex_chat(messages: list, max_tokens: int = 4000, temperature: float = 0.2,
                 req_timeout: int = 90) -> tuple[str, int]:
    api_key = os.environ.get("YANDEXGPT_API_KEY", "").strip()
    folder_id = os.environ.get("YANDEXGPT_FOLDER_ID", "").strip()
    if not api_key or not folder_id:
        raise RuntimeError("YANDEXGPT_API_KEY или YANDEXGPT_FOLDER_ID не заданы")
    yandex_messages = [{"role": m.get("role", "user"), "text": m.get("content", "")} for m in messages]
    payload = {
        "modelUri": f"gpt://{folder_id}/yandexgpt/latest",
        "completionOptions": {
            "stream": False,
            "temperature": temperature,
            "maxTokens": str(max_tokens),
        },
        "messages": yandex_messages,
    }
    req = urllib.request.Request(
        YANDEX_GPT_URL,
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        method="POST",
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Api-Key {api_key}",
            "x-folder-id": folder_id,
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=req_timeout) as r:
            body = json.loads(r.read().decode())
    except urllib.error.HTTPError as e:
        err_body = e.read().decode(errors="ignore")
        print(f"[YandexGPT] HTTP {e.code}: {err_body[:500]}")
        raise RuntimeError(f"YandexGPT HTTP {e.code}: {err_body[:300]}")
    alternatives = (body.get("result") or {}).get("alternatives") or []
    if not alternatives:
        raise RuntimeError(f"YandexGPT пустой ответ: {body}")
    text = alternatives[0].get("message", {}).get("text", "").strip()
    if not text:
        raise RuntimeError("YandexGPT вернул пустой текст")
    usage = (body.get("result") or {}).get("usage") or {}
    tokens_used = int(usage.get("totalTokens") or usage.get("completionTokens") or 0)
    return text, tokens_used


def get_gigachat_token():
    pass  # не используется, оставлен для совместимости ping-эндпоинта


def openrouter_chat(messages: list, max_tokens: int = 4000, temperature: float = 0.2,
                    req_timeout: int = 90) -> tuple[str, int]:
    return _yandex_chat(messages, max_tokens=max_tokens, temperature=temperature, req_timeout=req_timeout)


def gigachat_with_fallback(messages: list, max_tokens: int = 3000) -> tuple[str, int]:
    return _yandex_chat(messages, max_tokens=max_tokens)


def _repair_truncated_json(text: str) -> str:
    """Пытается закрыть обрезанный JSON: закрывает незакрытые строки/массивы/объекты."""
    # Убираем trailing мусор до последней «чистой» запятой или значения
    text = text.rstrip().rstrip(",").rstrip()
    # Считаем незакрытые скобки/кавычки
    in_string = False
    escape_next = False
    stack = []
    for ch in text:
        if escape_next:
            escape_next = False
            continue
        if ch == "\\" and in_string:
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch in ("{", "["):
            stack.append(ch)
        elif ch == "}":
            if stack and stack[-1] == "{":
                stack.pop()
        elif ch == "]":
            if stack and stack[-1] == "[":
                stack.pop()
    # Если внутри строки — закрываем её
    if in_string:
        text += '"'
    # Закрываем массивы и объекты в обратном порядке
    for ch in reversed(stack):
        text += "]" if ch == "[" else "}"
    return text


def extract_json(text: str) -> dict:
    text = text.strip()
    # Убираем markdown-блок если есть
    fence = re.search(r"```(?:json)?\s*(\{[\s\S]*)", text)
    if fence:
        text = fence.group(1)
        end_fence = text.find("```")
        if end_fence >= 0:
            text = text[:end_fence]
    # Берём от первой { до конца
    s = text.find("{")
    if s >= 0:
        text = text[s:]
    # Пробуем распарсить как есть
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # Пробуем закрыть обрезанный JSON
    fixed = _repair_truncated_json(text)
    return json.loads(fixed)


# ─── ПОИСК ИЗОБРАЖЕНИЙ ───────────────────────────────────────────────────────

def fetch_wikimedia(query: str, timeout: int = 5) -> bytes | None:
    """Ищет изображение через Wikimedia Commons."""
    try:
        # Попробуем сначала английский запрос (Wikimedia лучше по-английски)
        for q in [query]:
            search_q = urllib.parse.quote(q)
            url = (
                f"https://commons.wikimedia.org/w/api.php"
                f"?action=query&generator=search&gsrnamespace=6"
                f"&gsrsearch={search_q}&gsrlimit=5"
                f"&prop=imageinfo&iiprop=url|mime|size"
                f"&iiurlwidth=800&format=json&origin=*"
            )
            req = urllib.request.Request(url, headers={"User-Agent": "AOUSPT-Edu-Bot/1.0"})
            with urllib.request.urlopen(req, timeout=timeout) as r:
                data = json.loads(r.read().decode())

            pages = (data.get("query") or {}).get("pages") or {}
            candidates = []
            for page in pages.values():
                ii = (page.get("imageinfo") or [{}])[0]
                mime = ii.get("mime", "")
                thumb_url = ii.get("thumburl") or ii.get("url", "")
                size = ii.get("size", 0)
                if mime in ("image/jpeg", "image/png") and thumb_url and 8000 < size < 5_000_000:
                    candidates.append((size, thumb_url))

            # Берём самое «тяжёлое» (крупное) изображение
            candidates.sort(reverse=True)
            for _, thumb_url in candidates[:2]:
                try:
                    req2 = urllib.request.Request(thumb_url, headers={"User-Agent": "AOUSPT-Edu-Bot/1.0"})
                    with urllib.request.urlopen(req2, timeout=timeout) as r2:
                        return r2.read()
                except Exception:
                    continue
        return None
    except Exception:
        return None


def fetch_image_bytes(query: str, timeout: int = 5) -> bytes | None:
    """Ищет изображение: сначала Wikimedia, затем Wikipedia featured."""
    result = fetch_wikimedia(query, timeout)
    if result:
        return result
    # Fallback: ищем через Wikipedia API (статейные изображения)
    try:
        search_q = urllib.parse.quote(query)
        url = (
            f"https://ru.wikipedia.org/w/api.php"
            f"?action=query&titles={search_q}&prop=pageimages"
            f"&pithumbsize=700&format=json&origin=*"
        )
        req = urllib.request.Request(url, headers={"User-Agent": "AOUSPT-Edu-Bot/1.0"})
        with urllib.request.urlopen(req, timeout=timeout) as r:
            data = json.loads(r.read().decode())
        pages = (data.get("query") or {}).get("pages") or {}
        for page in pages.values():
            thumb = (page.get("thumbnail") or {}).get("source")
            if thumb:
                req2 = urllib.request.Request(thumb, headers={"User-Agent": "AOUSPT-Edu-Bot/1.0"})
                with urllib.request.urlopen(req2, timeout=timeout) as r2:
                    return r2.read()
    except Exception:
        pass
    return None


def fetch_images_parallel(slides: list, topic: str) -> list:
    """Параллельно загружает до 3 изображений для каждого слайда.
    Возвращает list[list[bytes]] — список фото для каждого слайда."""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    def get_imgs_for_slide(idx_slide):
        idx, slide = idx_slide
        queries = slide.get("image_queries") or []
        if not queries:
            old_q = slide.get("image_query") or slide.get("title", "")
            queries = [old_q, f"{topic} {slide.get('title', '')}", slide.get("title", "")]
        queries = [q for q in queries if q][:3]

        imgs = []
        seen_hashes = set()
        for q in queries:
            if len(imgs) >= 3:
                break
            img = fetch_image_bytes(q, timeout=4)
            if img:
                h = hash(img[:64])
                if h not in seen_hashes:
                    seen_hashes.add(h)
                    imgs.append(img)
        return idx, imgs

    results = [[] for _ in slides]
    with ThreadPoolExecutor(max_workers=min(len(slides) * 2, 12)) as ex:
        futures = {ex.submit(get_imgs_for_slide, (i, s)): i for i, s in enumerate(slides)}
        for future in as_completed(futures, timeout=20):
            try:
                idx, imgs = future.result()
                results[idx] = imgs
            except Exception:
                pass
    return results


def fetch_image_for_slide(slide_title: str, topic: str) -> bytes | None:
    """Совместимость со старым кодом."""
    return fetch_image_bytes(f"{slide_title} {topic}", timeout=5)


# ─── ГЕНЕРАЦИЯ СТРУКТУРЫ ─────────────────────────────────────────────────────

def generate_outline(topic: str, description: str, slides_count: int, audience: str) -> tuple[dict, int]:
    """Генерирует развёрнутую структуру презентации строго по ФГОС."""
    audience_str = audience or "школьники"

    system = (
        "Ты опытный методист и дизайнер учебных презентаций по ФГОС. "
        "Создаёшь насыщенный, информативный контент с конкретными фактами, датами, именами. "
        "Отвечаешь ТОЛЬКО валидным JSON без пояснений и markdown-блоков."
    )

    user = (
        f"Тема: {topic}. Аудитория: {audience_str}. Слайдов: {slides_count}.\n"
        + (f"Контекст учителя: {description[:800]}\n" if description else "")
        + 'Верни JSON строго в формате:\n'
        '{"subtitle":"развёрнутый подзаголовок 6-10 слов",'
        '"contents":["название раздела 1","название раздела 2","..."],'
        '"slides":[{'
        '"title":"заголовок слайда 4-7 слов",'
        '"bullets":["Развёрнутый тезис 1 — 15-25 слов с конкретным фактом","Тезис 2 — 15-25 слов","Тезис 3","Тезис 4","Тезис 5"],'
        '"fact":"Интересный факт или цитата эксперта — 15-25 слов",'
        '"image_queries":["конкретный запрос фото 1 на русском","конкретный запрос фото 2","конкретный запрос фото 3"]}],'
        '"conclusion":["Вывод 1 — 10-20 слов","Вывод 2","Вывод 3"]}\n'
        f"Ровно {slides_count} слайдов. bullets — ровно 5 штук на каждый слайд. image_queries — 3 разных запроса."
    )

    # ~300 токенов на слайд (5 тезисов × ~25 слов + заголовок + fact + 3 queries + запас)
    max_tok = min(300 * slides_count + 800, 7000)

    raw, tokens_used = gigachat_with_fallback(
        [{"role": "system", "content": system}, {"role": "user", "content": user}],
        max_tokens=max_tok,
    )

    try:
        data = extract_json(raw)
    except Exception as e:
        raise RuntimeError(f"Не удалось разобрать ответ ИИ как JSON: {e}. Ответ: {raw[:400]}")

    slides = data.get("slides") or []
    if not isinstance(slides, list) or not slides:
        raise RuntimeError("ИИ не вернул слайды")

    norm_slides = []
    for s in slides[:slides_count]:
        title = (s.get("title") or "").strip()
        bullets = s.get("bullets") or []
        if not isinstance(bullets, list):
            bullets = []
        bullets = [str(b).strip() for b in bullets if str(b).strip()][:7]
        fact = (s.get("fact") or "").strip()
        # Поддерживаем оба формата: новый image_queries (список) и старый image_query (строка)
        image_queries = s.get("image_queries") or []
        if not isinstance(image_queries, list) or not image_queries:
            single = (s.get("image_query") or title).strip()
            image_queries = [single, f"{topic} {title}", title] if single else [title, title, title]
        image_queries = [str(q).strip() for q in image_queries if str(q).strip()][:3]
        if not image_queries:
            image_queries = [title, topic, title]
        if title and bullets:
            norm_slides.append({
                "title": title,
                "bullets": bullets,
                "fact": fact,
                "image_queries": image_queries,
            })

    if not norm_slides:
        raise RuntimeError("ИИ вернул пустые слайды")

    contents = data.get("contents") or [s["title"] for s in norm_slides]

    return {
        "subtitle": (data.get("subtitle") or "").strip(),
        "contents": [str(c).strip() for c in contents if str(c).strip()][:slides_count],
        "slides": norm_slides,
        "conclusion": [str(c).strip() for c in (data.get("conclusion") or []) if str(c).strip()][:5],
    }, tokens_used


# ─── PPTX BUILDER ────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_solid_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_solid_fill(shape, rgb)
    return shape


def _add_oval(slide, x, y, w, h, rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    _set_solid_fill(shape, rgb)
    return shape


def _mix(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
    """Смешивает два цвета: t=0 → a, t=1 → b."""
    t = max(0.0, min(1.0, t))
    ar, ag, ab = a[0], a[1], a[2]
    br, bg, bb = b[0], b[1], b[2]
    return RGBColor(int(ar + (br - ar) * t),
                    int(ag + (bg - ag) * t),
                    int(ab + (bb - ab) * t))


# Палитра примитивов, из которых процедурно собирается уникальная композиция.
# Загружаем защищённо через getattr — состав enum может отличаться по версиям
# python-pptx, поэтому берём только реально доступные фигуры.
# Только простые геометрии, которые надёжно сериализуются во всех версиях
# python-pptx (без preset-guides, требующих доп. параметров).
_DECOR_PRIM_NAMES = {
    "oval":          "OVAL",
    "rect":          "RECTANGLE",
    "round_rect":    "ROUNDED_RECTANGLE",
    "parallelogram": "PARALLELOGRAM",
    "diamond":       "DIAMOND",
    "hexagon":       "HEXAGON",
    "trapezoid":     "TRAPEZOID",
    "pentagon":      "REGULAR_PENTAGON",
}
_DECOR_PRIMS = {}
for _k, _name in _DECOR_PRIM_NAMES.items():
    _shape = getattr(MSO_SHAPE, _name, None)
    if _shape is not None:
        _DECOR_PRIMS[_k] = _shape
# Гарантируем базовые фигуры
if not _DECOR_PRIMS:
    _DECOR_PRIMS = {"oval": MSO_SHAPE.OVAL, "rect": MSO_SHAPE.RECTANGLE}


def generate_decor_recipe(rnd: "random.Random") -> dict:
    """
    Процедурно генерирует УНИКАЛЬНУЮ графическую композицию во время запроса.
    Возвращает «рецепт» — список инструкций рисования в долях от размера слайда
    (0..1), цвета как роль+степень растворения. Каждый вызов даёт новый дизайн.
    """
    prims = list(_DECOR_PRIMS.keys())
    n = rnd.randint(3, 7)
    # Композиционная «зона тяготения» — графика концентрируется в случайном углу/крае
    anchors = ["tr", "tl", "br", "bl", "right", "left", "scatter"]
    anchor = rnd.choice(anchors)
    shapes = []
    for _ in range(n):
        prim = rnd.choice(prims)
        size = rnd.uniform(0.08, 0.55)        # доля от меньшей стороны
        # положение центра в зависимости от зоны тяготения
        if anchor == "tr":   cx, cy = rnd.uniform(0.6, 1.05), rnd.uniform(-0.05, 0.45)
        elif anchor == "tl": cx, cy = rnd.uniform(-0.05, 0.4), rnd.uniform(-0.05, 0.45)
        elif anchor == "br": cx, cy = rnd.uniform(0.6, 1.05), rnd.uniform(0.55, 1.05)
        elif anchor == "bl": cx, cy = rnd.uniform(-0.05, 0.4), rnd.uniform(0.55, 1.05)
        elif anchor == "right": cx, cy = rnd.uniform(0.7, 1.1), rnd.uniform(-0.1, 1.1)
        elif anchor == "left":  cx, cy = rnd.uniform(-0.1, 0.3), rnd.uniform(-0.1, 1.1)
        else:                cx, cy = rnd.uniform(-0.05, 1.05), rnd.uniform(-0.05, 1.05)
        shapes.append({
            "p": prim,
            "cx": round(cx, 4), "cy": round(cy, 4),
            "w": round(size, 4),
            "h": round(size * rnd.uniform(0.55, 1.6), 4),
            "rot": rnd.choice([0, 0, 15, 30, 45, -20, 60, 120]),
            "role": rnd.choice(["a2", "a2", "a3", "glow"]),
            "fade": round(rnd.uniform(0.5, 0.9), 3),  # насколько растворить в фоне
        })
    # Иногда добавляем тонкий «контур»-линию для современного штриха
    line = rnd.random() < 0.5
    return {"shapes": shapes, "anchor": anchor, "line": line,
            "line_rot": rnd.choice([0, 12, -12, 90])}


def _decorate(slide, theme: dict, on_dark: bool):
    """
    Рисует процедурно сгенерированную композицию из theme['decor'] (рецепт).
    on_dark — тёмный фон (титул) или светлый (контент): подбираем мягкость.
    """
    recipe = theme.get("decor")
    if not recipe or not recipe.get("shapes"):
        return
    base = theme["title_bg"] if on_dark else theme["bg"]
    roles = {
        "a2":   theme["accent2"],
        "a3":   theme.get("accent3", theme["accent2"]),
        "glow": theme.get("stripe", theme["accent2"]),
    }
    side = min(int(SLIDE_W), int(SLIDE_H))

    for sp in recipe["shapes"]:
        prim = _DECOR_PRIMS.get(sp.get("p"), MSO_SHAPE.OVAL)
        w = int(sp["w"] * side)
        h = int(sp["h"] * side)
        x = int(sp["cx"] * int(SLIDE_W) - w / 2)
        y = int(sp["cy"] * int(SLIDE_H) - h / 2)
        col = roles.get(sp.get("role"), theme["accent2"])
        # На светлом фоне растворяем сильнее, чтобы графика не перебивала текст
        fade = sp.get("fade", 0.7)
        fade = min(0.92, fade + (0.12 if not on_dark else 0.0))
        soft = _mix(col, base, fade)
        try:
            shape = slide.shapes.add_shape(prim, x, y, max(w, 1), max(h, 1))
            _set_solid_fill(shape, soft)
            if sp.get("rot"):
                shape.rotation = sp["rot"]
        except Exception:
            continue

    if recipe.get("line"):
        try:
            col = _mix(theme["accent2"], base, 0.4 if on_dark else 0.5)
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        int(SLIDE_W * 0.62), int(SLIDE_H * 0.1),
                                        Emu(12000), int(SLIDE_H * 0.8))
            _set_solid_fill(ln, col)
            if recipe.get("line_rot"):
                ln.rotation = recipe["line_rot"]
        except Exception:
            pass


def _add_text(slide, x, y, w, h, text: str, *, size: int = 18, bold: bool = False,
              color: RGBColor = None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font: str = "Calibri", italic: bool = False):
    if color is None:
        color = RGBColor(0x22, 0x2A, 0x35)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_bullets(slide, x, y, w, h, bullets: list, *, size: int = 17,
                 color: RGBColor = None, font: str = "Calibri",
                 accent2: RGBColor = None, space_after: int = 8):
    if color is None:
        color = RGBColor(0x22, 0x2A, 0x35)
    if accent2 is None:
        accent2 = RGBColor(0xC2, 0x8B, 0x42)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.name = font
        dot.font.size = Pt(size)
        dot.font.color.rgb = accent2
        dot.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _add_image_to_slide(slide, img_bytes: bytes, x, y, w, h):
    """Добавляет изображение на слайд из bytes."""
    try:
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, x, y, w, h)
        return True
    except Exception:
        return False


def _add_image_placeholder(slide, x, y, w, h, theme: dict, label: str = ""):
    """Заглушка если фото не найдено — цветной прямоугольник."""
    _add_rect(slide, x, y, w, h, theme["accent"])
    if label:
        _add_text(slide, x + Inches(0.1), y + h // 2 - Inches(0.2),
                  w - Inches(0.2), Inches(0.4),
                  label, size=11, color=theme["title_sub"],
                  align=PP_ALIGN.CENTER, italic=True)


def _footer(slide, teacher_name: str, teacher_school: str, theme: dict):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), SLIDE_H - Inches(0.52),
        Inches(12.33), Emu(7620)
    )
    _set_solid_fill(line, theme["accent2"])
    parts = [p for p in [teacher_name, teacher_school] if p]
    text = "  ·  ".join(parts)
    if text:
        _add_text(slide,
                  Inches(0.5), SLIDE_H - Inches(0.44),
                  Inches(12.33), Inches(0.35),
                  text, size=11, color=theme["muted"], align=PP_ALIGN.RIGHT)


def _slide_header(slide, title: str, slide_num: int, total: int,
                  teacher_name: str, theme: dict):
    """Шапка слайда: стиль зависит от layout темы."""
    layout = theme.get("layout", "top_banner")

    if layout == "sidebar_dark":
        # Тёмный вертикальный сайдбар слева с большим номером
        sidebar_w = Inches(1.8)
        _add_rect(slide, 0, 0, sidebar_w, SLIDE_H, theme["accent"])
        _add_rect(slide, sidebar_w, 0, Inches(0.06), SLIDE_H, theme["accent2"])
        _add_text(slide, 0, Inches(0.3), sidebar_w, Inches(1.0),
                  f"{slide_num:02d}", size=52, bold=True, color=theme["accent2"],
                  align=PP_ALIGN.CENTER)
        _add_text(slide, 0, Inches(1.3), sidebar_w, Inches(0.5),
                  f"/ {total:02d}", size=14, color=theme["title_sub"],
                  align=PP_ALIGN.CENTER)
        # Заголовок в правой части
        _add_rect(slide, sidebar_w + Inches(0.06), 0,
                  SLIDE_W - sidebar_w - Inches(0.06), Inches(1.2), theme["accent"])
        _add_text(slide, sidebar_w + Inches(0.2), Inches(0.15),
                  SLIDE_W - sidebar_w - Inches(0.4), Inches(0.9),
                  title, size=24, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        if teacher_name:
            _add_text(slide, sidebar_w + Inches(0.2), Inches(0.15),
                      SLIDE_W - sidebar_w - Inches(0.4), Inches(0.9),
                      teacher_name, size=10, color=theme["title_sub"],
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

    elif layout == "split_diagonal":
        # Горизонтальная шапка с диагональным акцентом
        _add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), theme["accent"])
        # Диагональный декор-прямоугольник
        _add_rect(slide, Inches(10.5), 0, Inches(2.83), Inches(1.5), theme["accent2"])
        _add_text(slide, Inches(10.55), 0, Inches(2.7), Inches(1.5),
                  f"{slide_num:02d}/{total:02d}", size=28, bold=True,
                  color=theme["accent"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, Inches(0.4), Inches(0.15), Inches(10.0), Inches(1.1),
                  title, size=26, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, 0, Inches(1.5), Inches(0.14), SLIDE_H - Inches(1.5), theme["accent2"])

    elif layout == "left_header":
        # Широкая полоса с крупным цветным номером в кружке
        _add_rect(slide, 0, 0, SLIDE_W, Inches(1.55), theme["accent"])
        # Цветной блок-нумератор
        _add_rect(slide, Inches(0.25), Inches(0.22), Inches(0.95), Inches(0.95), theme["accent2"])
        _add_text(slide, Inches(0.25), Inches(0.22), Inches(0.95), Inches(0.95),
                  str(slide_num), size=24, bold=True, color=theme["accent"],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, Inches(1.35), Inches(0.22), Inches(10.3), Inches(1.05),
                  title, size=26, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, Inches(1.35), Inches(0.22), Inches(11.6), Inches(1.05),
                  f"{total}", size=11, color=theme["title_sub"],
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
        _add_rect(slide, 0, Inches(1.55), SLIDE_W, Emu(9144), theme["accent2"])

    else:  # top_banner (default)
        _add_rect(slide, 0, 0, SLIDE_W, Inches(1.45), theme["accent"])
        _add_rect(slide, 0, Inches(1.45), Inches(0.18),
                  SLIDE_H - Inches(1.45), theme["accent2"])
        _add_text(slide, Inches(0.3), Inches(0.28), Inches(2), Inches(0.5),
                  f"{slide_num:02d} / {total:02d}",
                  size=12, bold=True, color=theme["accent2"], align=PP_ALIGN.LEFT)
        _add_text(slide, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.9),
                  title, size=26, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        if teacher_name:
            _add_text(slide, Inches(0.3), Inches(0.28), Inches(12.73), Inches(0.5),
                      teacher_name, size=11, color=theme["title_sub"], align=PP_ALIGN.RIGHT)


def _content_area_top(theme: dict) -> float:
    """Возвращает Y-начало контентной зоны в зависимости от layout."""
    layout = theme.get("layout", "top_banner")
    if layout == "sidebar_dark":
        return Inches(0.2)       # контент начинается сразу, сайдбар слева
    elif layout == "split_diagonal":
        return Inches(1.6)
    elif layout == "left_header":
        return Inches(1.7)
    return Inches(1.58)          # top_banner


def _content_area_x(theme: dict) -> float:
    layout = theme.get("layout", "top_banner")
    if layout == "sidebar_dark":
        return Inches(2.0)       # после сайдбара
    return Inches(0.32)


def _content_area_w(theme: dict, has_photo: bool) -> float:
    layout = theme.get("layout", "top_banner")
    sidebar = layout == "sidebar_dark"
    base_x = Inches(2.0) if sidebar else Inches(0.32)
    total_w = SLIDE_W - base_x - Inches(0.2)
    if has_photo:
        return total_w - Inches(4.7)
    return total_w


def _img_framed(slide, img_bytes, x, y, w, h, theme, rnd=None):
    """Фото с рамкой-акцентом. Стиль рамки слегка варьируется."""
    style = rnd.choice(["full", "corner", "thin", "none"]) if rnd else "full"
    if style == "full":
        pad = Emu(55000)
        _add_rect(slide, x - pad, y - pad, w + pad * 2, h + pad * 2, theme["accent2"])
    elif style == "corner":
        _add_rect(slide, x - Emu(60000), y - Emu(60000), Inches(0.5), Inches(0.5), theme["accent2"])
        _add_rect(slide, x + w - Inches(0.5) + Emu(60000), y + h - Inches(0.5) + Emu(60000),
                  Inches(0.5), Inches(0.5), theme["accent2"])
    elif style == "thin":
        _add_rect(slide, x - Emu(25000), y - Emu(25000), w + Emu(50000), h + Emu(50000), theme["accent2"])
    _add_image_to_slide(slide, img_bytes, x, y, w, h)


def _fact_card(slide, x, y, w, h, fact, theme, *, size=12, layout_style=None, rnd=None):
    """Карточка «интересный факт» — вид варьируется (полоса сверху / сбоку / звезда)."""
    if h <= Inches(0.35) or not fact:
        return
    style = layout_style or (rnd.choice(["top_bar", "side_bar", "star", "quote"]) if rnd else "top_bar")
    _add_rect(slide, x, y, w, h, theme["accent"])
    if style == "top_bar":
        _add_rect(slide, x, y, w, Emu(120000), theme["accent2"])
        _add_text(slide, x + Inches(0.15), y + Emu(12000), w - Inches(0.3), Emu(105000),
                  "ИНТЕРЕСНЫЙ ФАКТ", size=9, bold=True, color=theme["accent2"])
        _add_text(slide, x + Inches(0.15), y + Inches(0.3), w - Inches(0.3), h - Inches(0.38),
                  fact, size=size, color=theme["white"], italic=True, anchor=MSO_ANCHOR.TOP)
    elif style == "side_bar":
        _add_rect(slide, x, y, Emu(90000), h, theme["accent2"])
        _add_text(slide, x + Inches(0.22), y + Inches(0.12), w - Inches(0.35), h - Inches(0.24),
                  fact, size=size, color=theme["white"], italic=True, anchor=MSO_ANCHOR.MIDDLE)
    elif style == "quote":
        _add_text(slide, x + Inches(0.12), y + Emu(10000), Inches(0.8), h,
                  "❝", size=26, bold=True, color=theme["accent2"], anchor=MSO_ANCHOR.TOP)
        _add_text(slide, x + Inches(0.7), y + Inches(0.15), w - Inches(0.85), h - Inches(0.3),
                  fact, size=size, color=theme["white"], italic=True, anchor=MSO_ANCHOR.MIDDLE)
    else:  # star
        _add_text(slide, x + Inches(0.12), y, Inches(0.9), h, "★", size=18, bold=True,
                  color=theme["accent2"], anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, x + Inches(1.0), y + Emu(20000), w - Inches(1.15), h - Emu(40000),
                  fact, size=size, color=theme["white"], italic=True, anchor=MSO_ANCHOR.MIDDLE)


def place_slide_content(slide, theme, rnd, ct, cx, content_h, bullets, fact, imgs):
    """
    Процедурно размещает контент слайда: расположение фото, текста и факта
    выбирается СЛУЧАЙНО для каждого слайда. Один и тот же слайд не повторяет
    композицию соседнего — нет единого формата.
    """
    right_edge = SLIDE_W - Inches(0.25)
    avail_w = right_edge - cx
    n = len(imgs)

    if n == 0:
        # Варианты без фото: факт справа / факт снизу-полосой / только текст по центру
        modes = ["fact_right", "fact_bottom", "text_full", "fact_left"]
        mode = rnd.choice(modes) if fact else "text_full"
        if mode == "fact_right":
            fw = Inches(rnd.uniform(3.6, 4.3))
            fx = right_edge - fw
            fh = Inches(rnd.uniform(2.6, 3.4))
            _fact_card(slide, fx, ct + Inches(0.1), fw, fh, fact, theme, size=13, rnd=rnd)
            _add_bullets(slide, cx, ct, fx - cx - Inches(0.3), content_h, bullets,
                         size=16, color=theme["text"], accent2=theme["accent2"], space_after=10)
        elif mode == "fact_left":
            fw = Inches(rnd.uniform(3.4, 4.0))
            fh = Inches(rnd.uniform(2.6, 3.3))
            _fact_card(slide, cx, ct + Inches(0.1), fw, fh, fact, theme, size=13, rnd=rnd)
            tx = cx + fw + Inches(0.3)
            _add_bullets(slide, tx, ct, right_edge - tx, content_h, bullets,
                         size=16, color=theme["text"], accent2=theme["accent2"], space_after=10)
        elif mode == "fact_bottom":
            th = content_h * rnd.uniform(0.55, 0.66)
            _add_bullets(slide, cx, ct, avail_w, th, bullets,
                         size=16, color=theme["text"], accent2=theme["accent2"], space_after=9)
            fy = ct + th + Inches(0.15)
            _fact_card(slide, cx, fy, avail_w, SLIDE_H - fy - Inches(0.55), fact, theme, size=13, rnd=rnd)
        else:
            _add_bullets(slide, cx, ct, avail_w, content_h, bullets,
                         size=17, color=theme["text"], accent2=theme["accent2"], space_after=11)
        return

    if n == 1:
        # Один снимок: слева/справа/сверху/крупный-фон — выбираем случайно
        modes = ["photo_right", "photo_left", "photo_top", "photo_hero"]
        mode = rnd.choice(modes)
        if mode in ("photo_right", "photo_left"):
            pw = Inches(rnd.uniform(4.6, 5.6))
            ph = min(content_h - Inches(0.1), Inches(rnd.uniform(4.2, 5.0)))
            py = ct + Emu(rnd.randint(0, 40000))
            if mode == "photo_right":
                px = right_edge - pw
                tx, tw = cx, px - cx - Inches(0.25)
            else:
                px = cx
                tx, tw = cx + pw + Inches(0.3), right_edge - (cx + pw + Inches(0.3))
            _img_framed(slide, imgs[0], px, py, pw, ph, theme, rnd)
            bh = Inches(3.7) if fact else content_h
            _add_bullets(slide, tx, ct, max(tw, Inches(3.2)), bh, bullets,
                         size=15, color=theme["text"], accent2=theme["accent2"], space_after=9)
            if fact:
                fy = ct + Inches(3.85)
                _fact_card(slide, tx, fy, max(tw, Inches(3.2)), SLIDE_H - fy - Inches(0.55),
                           fact, theme, size=12, rnd=rnd)
        elif mode == "photo_top":
            ph = content_h * rnd.uniform(0.42, 0.5)
            _img_framed(slide, imgs[0], cx, ct, avail_w, ph, theme, rnd)
            ty = ct + ph + Inches(0.25)
            _add_bullets(slide, cx, ty, avail_w, SLIDE_H - ty - Inches(0.55), bullets,
                         size=15, color=theme["text"], accent2=theme["accent2"], space_after=8)
        else:  # photo_hero — фото крупным блоком слева, текст узкой колонкой
            pw = Inches(rnd.uniform(6.5, 7.5))
            _img_framed(slide, imgs[0], cx, ct, pw, content_h, theme, rnd)
            tx = cx + pw + Inches(0.3)
            tw = right_edge - tx
            _add_bullets(slide, tx, ct, tw, content_h, bullets,
                         size=14, color=theme["text"], accent2=theme["accent2"], space_after=8)
        return

    if n == 2:
        modes = ["col_right", "col_left", "top_pair", "diagonal"]
        mode = rnd.choice(modes)
        if mode in ("col_right", "col_left"):
            colw = Inches(rnd.uniform(3.9, 4.6))
            ph = (content_h - Inches(0.15)) / 2
            if mode == "col_right":
                colx = right_edge - colw
                tx, tw = cx, colx - cx - Inches(0.3)
            else:
                colx = cx
                tx, tw = cx + colw + Inches(0.3), right_edge - (cx + colw + Inches(0.3))
            _img_framed(slide, imgs[0], colx, ct, colw, ph, theme, rnd)
            _img_framed(slide, imgs[1], colx, ct + ph + Inches(0.15), colw, ph, theme, rnd)
            bh = Inches(3.6) if fact else content_h
            _add_bullets(slide, tx, ct, max(tw, Inches(3.0)), bh, bullets,
                         size=15, color=theme["text"], accent2=theme["accent2"], space_after=9)
            if fact:
                fy = ct + Inches(3.75)
                _fact_card(slide, tx, fy, max(tw, Inches(3.0)), SLIDE_H - fy - Inches(0.55),
                           fact, theme, size=12, rnd=rnd)
        elif mode == "top_pair":
            ph = content_h * rnd.uniform(0.4, 0.47)
            half = (avail_w - Inches(0.2)) / 2
            _img_framed(slide, imgs[0], cx, ct, half, ph, theme, rnd)
            _img_framed(slide, imgs[1], cx + half + Inches(0.2), ct, half, ph, theme, rnd)
            ty = ct + ph + Inches(0.25)
            _add_bullets(slide, cx, ty, avail_w, SLIDE_H - ty - Inches(0.55), bullets,
                         size=15, color=theme["text"], accent2=theme["accent2"], space_after=8)
        else:  # diagonal — одно фото сверху-слева, второе снизу-справа, текст обтекает
            ph = content_h * 0.46
            pw = avail_w * 0.5
            _img_framed(slide, imgs[0], cx, ct, pw, ph, theme, rnd)
            _img_framed(slide, imgs[1], right_edge - pw, SLIDE_H - Inches(0.55) - ph, pw, ph, theme, rnd)
            _add_bullets(slide, cx + pw + Inches(0.25), ct, right_edge - (cx + pw + Inches(0.25)),
                         ph, bullets[:3], size=14, color=theme["text"],
                         accent2=theme["accent2"], space_after=7)
            _add_bullets(slide, cx, ct + ph + Inches(0.2), pw + Inches(0.5),
                         SLIDE_H - (ct + ph + Inches(0.2)) - Inches(0.55), bullets[3:],
                         size=14, color=theme["text"], accent2=theme["accent2"], space_after=7)
        return

    # n >= 3
    modes = ["strip_bottom", "strip_top", "grid_right", "mosaic"]
    mode = rnd.choice(modes)
    three = imgs[:3]
    if mode in ("strip_bottom", "strip_top"):
        strip_h = content_h * rnd.uniform(0.4, 0.48)
        pw = (avail_w - Inches(0.24)) / 3
        if mode == "strip_bottom":
            strip_y = SLIDE_H - Inches(0.55) - strip_h
            _add_bullets(slide, cx, ct, avail_w, strip_y - ct - Inches(0.15), bullets,
                         size=14, color=theme["text"], accent2=theme["accent2"], space_after=7)
        else:
            strip_y = ct
            ty = ct + strip_h + Inches(0.2)
            _add_bullets(slide, cx, ty, avail_w, SLIDE_H - ty - Inches(0.55), bullets,
                         size=14, color=theme["text"], accent2=theme["accent2"], space_after=7)
        for pi, img in enumerate(three):
            _img_framed(slide, img, cx + pi * (pw + Inches(0.12)), strip_y, pw, strip_h, theme, rnd)
    elif mode == "grid_right":
        colw = Inches(rnd.uniform(3.8, 4.4))
        colx = right_edge - colw
        ph = (content_h - Inches(0.24)) / 3
        for pi, img in enumerate(three):
            _img_framed(slide, img, colx, ct + pi * (ph + Inches(0.12)), colw, ph, theme, rnd)
        _add_bullets(slide, cx, ct, colx - cx - Inches(0.3), content_h, bullets,
                     size=15, color=theme["text"], accent2=theme["accent2"], space_after=9)
    else:  # mosaic — крупное фото + два мелких, текст в свободной зоне
        big_w = avail_w * 0.5
        big_h = content_h * 0.62
        _img_framed(slide, three[0], cx, ct, big_w, big_h, theme, rnd)
        sw = big_w
        sh = (content_h - big_h - Inches(0.15))
        _img_framed(slide, three[1], cx, ct + big_h + Inches(0.15), (sw - Inches(0.1)) / 2, sh, theme, rnd)
        _img_framed(slide, three[2], cx + (sw - Inches(0.1)) / 2 + Inches(0.1),
                    ct + big_h + Inches(0.15), (sw - Inches(0.1)) / 2, sh, theme, rnd)
        tx = cx + big_w + Inches(0.3)
        _add_bullets(slide, tx, ct, right_edge - tx, content_h, bullets,
                     size=14, color=theme["text"], accent2=theme["accent2"], space_after=8)


def build_pptx(topic: str, subtitle: str, contents: list, slides_data: list,
               conclusion: list, teacher_name: str, teacher_school: str,
               theme: dict, images: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total_content = len(slides_data)
    layout = theme.get("layout", "top_banner")
    parts = [p for p in [teacher_name, teacher_school] if p]
    footer_text = "   ·   ".join(parts)

    # ── 1. Титульный слайд ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["title_bg"])
    _decorate(slide, theme, on_dark=True)

    if layout == "sidebar_dark":
        # Вертикальный акцент справа
        _add_rect(slide, SLIDE_W - Inches(0.55), 0, Inches(0.55), SLIDE_H, theme["accent2"])
        _add_rect(slide, 0, Inches(3.5), SLIDE_W - Inches(0.55), Inches(0.06), theme["accent2"])
        _add_text(slide, Inches(0.7), Inches(0.9), Inches(9), Inches(0.45),
                  theme["label"], size=11, bold=True, color=theme["accent2"])
        _add_text(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(2.1),
                  topic, size=44, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        if subtitle:
            _add_text(slide, Inches(0.7), Inches(3.65), Inches(11.5), Inches(0.9),
                      subtitle, size=17, color=theme["title_sub"])

    elif layout == "split_diagonal":
        # Нижний блок с диагональным стыком — визуально через два прямоугольника
        _add_rect(slide, 0, Inches(4.0), SLIDE_W, Inches(3.5), theme["accent2"])
        _add_rect(slide, 0, Inches(4.0), Inches(7.0), Inches(3.5), theme["accent"])
        _add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, theme["accent2"])
        _add_text(slide, Inches(0.7), Inches(0.8), Inches(9), Inches(0.5),
                  theme["label"], size=12, bold=True, color=theme["accent2"])
        _add_text(slide, Inches(0.7), Inches(1.4), Inches(11.8), Inches(2.5),
                  topic, size=44, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        if subtitle:
            _add_text(slide, Inches(0.7), Inches(4.1), Inches(6.5), Inches(0.9),
                      subtitle, size=16, color=theme["white"])

    elif layout == "left_header":
        # Крупный цветной блок слева
        _add_rect(slide, 0, 0, Inches(5.5), SLIDE_H, theme["accent"])
        _add_rect(slide, Inches(5.5), 0, Inches(0.08), SLIDE_H, theme["accent2"])
        _add_text(slide, Inches(0.4), Inches(1.0), Inches(5.0), Inches(0.5),
                  theme["label"], size=11, bold=True, color=theme["accent2"])
        _add_text(slide, Inches(0.4), Inches(1.6), Inches(4.8), Inches(3.5),
                  topic, size=36, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        if subtitle:
            _add_text(slide, Inches(0.4), Inches(5.5), Inches(4.8), Inches(0.9),
                      subtitle, size=15, color=theme["title_sub"])
        # Правая часть — дата и учитель
        _add_text(slide, Inches(5.8), Inches(2.5), Inches(7.0), Inches(1.5),
                  datetime.now().strftime("%d.%m.%Y"), size=28,
                  color=theme["title_sub"], align=PP_ALIGN.LEFT)

    else:  # top_banner
        _add_rect(slide, 0, Inches(4.6), SLIDE_W, Inches(2.9), theme["accent"])
        _add_rect(slide, 0, Inches(4.55), SLIDE_W, Inches(0.08), theme["accent2"])
        _add_rect(slide, 0, 0, Inches(0.45), SLIDE_H, theme["accent2"])
        _add_text(slide, Inches(0.7), Inches(1.55), Inches(10), Inches(0.4),
                  theme["label"], size=11, bold=True, color=theme["accent2"])
        _add_text(slide, Inches(0.7), Inches(2.05), Inches(11.8), Inches(2.2),
                  topic, size=42, bold=True, color=theme["white"],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        if subtitle:
            _add_text(slide, Inches(0.7), Inches(4.75), Inches(11.8), Inches(0.8),
                      subtitle, size=17, color=theme["title_sub"])

    if footer_text and layout != "left_header":
        _add_text(slide, Inches(0.7), SLIDE_H - Inches(0.75),
                  Inches(10), Inches(0.45),
                  footer_text, size=13, color=theme["title_sub"])
        _add_text(slide, Inches(0.7), SLIDE_H - Inches(0.75),
                  Inches(12.13), Inches(0.45),
                  datetime.now().strftime("%d.%m.%Y"),
                  size=12, color=theme["title_sub"], align=PP_ALIGN.RIGHT)

    # ── 2. Слайд «Содержание» ────────────────────────────────────────────
    if contents:
        slide = prs.slides.add_slide(blank)
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg"])

        if layout == "sidebar_dark":
            _add_rect(slide, 0, 0, Inches(1.8), SLIDE_H, theme["accent"])
            _add_rect(slide, Inches(1.8), 0, Inches(0.06), SLIDE_H, theme["accent2"])
            _add_text(slide, 0, Inches(0.5), Inches(1.8), Inches(1.2),
                      "СО\nДЕР\nЖА\nНИЕ", size=12, bold=True, color=theme["accent2"],
                      align=PP_ALIGN.CENTER)
            cx = Inches(2.1)
            cw = Inches(10.8)
            cy = Inches(0.4)
        else:
            _add_rect(slide, 0, 0, SLIDE_W, Inches(1.45), theme["accent"])
            _add_rect(slide, 0, Inches(1.45), Inches(0.18),
                      SLIDE_H - Inches(1.45), theme["accent2"])
            _add_text(slide, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.9),
                      "Содержание урока", size=28, bold=True, color=theme["white"],
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            cx = Inches(0.38)
            cw = Inches(6.0)
            cy = Inches(1.65)

        mid = (len(contents) + 1) // 2
        left_items = contents[:mid]
        right_items = contents[mid:]

        if layout == "sidebar_dark":
            # Нумерованный список крупными карточками
            for i, item in enumerate(contents):
                card_y = cy + i * Inches(0.75)
                if card_y + Inches(0.65) > SLIDE_H - Inches(0.3):
                    break
                _add_rect(slide, cx, card_y, cw, Inches(0.62), theme["card_bg"])
                _add_text(slide, cx + Inches(0.1), card_y, Inches(0.6), Inches(0.62),
                          f"{i+1:02d}", size=14, bold=True, color=theme["accent2"],
                          anchor=MSO_ANCHOR.MIDDLE)
                _add_text(slide, cx + Inches(0.75), card_y, cw - Inches(0.85), Inches(0.62),
                          item, size=15, color=theme["text"], anchor=MSO_ANCHOR.MIDDLE)
        else:
            _add_bullets(slide, cx, cy, cw, Inches(5.4),
                         left_items, size=17, color=theme["text"],
                         accent2=theme["accent2"], space_after=14)
            if right_items:
                _add_bullets(slide, Inches(6.8), cy, cw, Inches(5.4),
                             right_items, size=17, color=theme["text"],
                             accent2=theme["accent2"], space_after=14)
        _footer(slide, teacher_name, teacher_school, theme)

    # ── 3. Содержательные слайды (современный layout) ────────────────────
    for idx, s in enumerate(slides_data, start=1):
        slide = prs.slides.add_slide(blank)
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg"])
        _decorate(slide, theme, on_dark=False)
        _slide_header(slide, s["title"], idx, total_content, teacher_name, theme)

        slide_imgs = images.get(idx) or []  # list[bytes], до 3 штук
        fact = s.get("fact", "")
        bullets = s["bullets"]
        ct = _content_area_top(theme)
        cx = _content_area_x(theme)
        content_h = SLIDE_H - ct - Inches(0.6)

        # ── Процедурный per-slide макет: расположение фото, текста и факта
        #    выбирается СЛУЧАЙНО для каждого слайда — единого формата нет. ──
        slide_rnd = random.Random(f"{theme.get('name','')}|{s['title']}|{idx}|{random.random()}")
        place_slide_content(slide, theme, slide_rnd, ct, cx, content_h, bullets, fact, slide_imgs)

        _footer(slide, teacher_name, teacher_school, theme)

    # ── 4. Слайд выводов ─────────────────────────────────────────────────
    if conclusion:
        slide = prs.slides.add_slide(blank)
        _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg"])

        if layout == "sidebar_dark":
            _add_rect(slide, 0, 0, Inches(1.8), SLIDE_H, theme["accent"])
            _add_rect(slide, Inches(1.8), 0, Inches(0.06), SLIDE_H, theme["accent2"])
            _add_text(slide, 0, Inches(0.5), Inches(1.8), Inches(2.0),
                      "ИТО\nГИ", size=20, bold=True, color=theme["accent2"],
                      align=PP_ALIGN.CENTER)
            # Крупные карточки выводов
            for i, item in enumerate(conclusion[:4]):
                cy = Inches(0.4) + i * Inches(1.65)
                _add_rect(slide, Inches(2.1), cy, Inches(10.8), Inches(1.5), theme["card_bg"])
                _add_rect(slide, Inches(2.1), cy, Inches(0.5), Inches(1.5), theme["accent2"])
                _add_text(slide, Inches(2.65), cy + Emu(50000), Inches(10.0),
                          Inches(1.4), item, size=16, color=theme["text"],
                          anchor=MSO_ANCHOR.MIDDLE)
        elif layout in ("split_diagonal", "left_header"):
            # Двухколонная сетка карточек
            _add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), theme["accent"])
            _add_rect(slide, 0, Inches(1.3), SLIDE_W, Emu(9144), theme["accent2"])
            _add_text(slide, Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.9),
                      "Ключевые выводы урока", size=28, bold=True, color=theme["white"],
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            cols = [(Inches(0.3), Inches(6.3)), (Inches(6.8), Inches(6.3))]
            for i, item in enumerate(conclusion[:4]):
                col_x, col_w = cols[i % 2]
                row_y = Inches(1.55) + (i // 2) * Inches(2.8)
                _add_rect(slide, col_x, row_y, col_w, Inches(2.5), theme["card_bg"])
                _add_rect(slide, col_x, row_y, col_w, Emu(110000), theme["accent2"])
                _add_text(slide, col_x + Inches(0.15), row_y + Inches(0.2),
                          col_w - Inches(0.3), Inches(2.2),
                          item, size=15, color=theme["text"], anchor=MSO_ANCHOR.TOP)
        else:
            _add_rect(slide, 0, 0, SLIDE_W, Inches(1.45), theme["accent"])
            _add_rect(slide, 0, Inches(1.45), Inches(0.18),
                      SLIDE_H - Inches(1.45), theme["accent2"])
            _add_text(slide, Inches(0.3), Inches(0.28), Inches(3), Inches(0.5),
                      "ИТОГИ УРОКА", size=11, bold=True, color=theme["accent2"])
            _add_text(slide, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.9),
                      "Ключевые выводы", size=28, bold=True, color=theme["white"],
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            _add_bullets(slide, Inches(0.32), Inches(1.65),
                         Inches(12.6), Inches(5.1),
                         conclusion, size=19, color=theme["text"],
                         accent2=theme["accent2"], space_after=12)

        _footer(slide, teacher_name, teacher_school, theme)

    # ── 5. Финальный слайд ───────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["title_bg"])

    if layout == "sidebar_dark":
        _add_rect(slide, SLIDE_W - Inches(0.55), 0, Inches(0.55), SLIDE_H, theme["accent2"])
        _add_rect(slide, 0, Inches(3.5), SLIDE_W - Inches(0.55), Inches(0.06), theme["accent2"])
    elif layout == "split_diagonal":
        _add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, theme["accent2"])
        _add_rect(slide, 0, Inches(4.0), SLIDE_W, Inches(0.06), theme["accent2"])
    elif layout == "left_header":
        _add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, theme["accent2"])
        _add_rect(slide, 0, SLIDE_H - Inches(2.0), SLIDE_W, Inches(0.06), theme["accent2"])
    else:
        _add_rect(slide, 0, Inches(4.55), SLIDE_W, Inches(0.08), theme["accent2"])
        _add_rect(slide, 0, 0, Inches(0.45), SLIDE_H, theme["accent2"])

    _add_text(slide, Inches(0.7), Inches(2.5), Inches(11.8), Inches(1.6),
              "Спасибо за внимание!", size=52, bold=True,
              color=theme["white"], align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(0.7), Inches(4.2), Inches(11.8), Inches(0.7),
              "Вопросы и обсуждение", size=22,
              color=theme["title_sub"], align=PP_ALIGN.CENTER)
    if footer_text:
        _add_text(slide, Inches(0.7), SLIDE_H - Inches(0.75),
                  Inches(12.13), Inches(0.45),
                  footer_text, size=13, color=theme["title_sub"], align=PP_ALIGN.CENTER)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def safe_filename(s: str, max_len: int = 60) -> str:
    s = re.sub(r"[\\/:*?\"<>|]+", " ", s).strip()
    s = re.sub(r"\s+", " ", s)
    return s[:max_len] or "Презентация"


# ─── HANDLER ─────────────────────────────────────────────────────────────────

def _parse_body(event: dict) -> dict:
    raw = event.get("body") or ""
    if event.get("isBase64Encoded"):
        try:
            raw = base64.b64decode(raw).decode()
        except Exception:
            raw = ""
    try:
        body = json.loads(raw) if raw else {}
        if isinstance(body, str):
            body = json.loads(body)
        return body
    except Exception:
        return {}


def handler(event: dict, context) -> dict:
    """
    Генерация PPTX разбита на 2 шага, чтобы уложиться в таймаут платформы:
    POST ?action=outline  — запрос к GigaChat, возвращает JSON-структуру (~60 сек)
    POST ?action=build    — скачивает фото и собирает PPTX из готовой структуры (~15 сек)
    POST (без action)     — legacy: оба шага вместе (для совместимости)
    """
    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": CORS, "body": ""}

    method = event.get("httpMethod", "POST")
    qs = event.get("queryStringParameters") or {}
    action = (qs.get("action") or "").strip().lower()

    if action in ("ping", "warmup"):
        try:
            api_key = os.environ.get("YANDEXGPT_API_KEY", "").strip()
            if not api_key:
                raise RuntimeError("YANDEXGPT_API_KEY не задан")
            return _resp(200, {"ok": True, "service": "YandexGPT"})
        except Exception as e:
            return _resp(500, {"ok": False, "error": str(e)})

    if method != "POST":
        return _resp(405, {"error": "Метод не поддерживается"})

    body = _parse_body(event)

    # ── ШАГ 1: outline — только GigaChat (~60 сек) ──────────────────────────
    if action == "outline":
        login = (body.get("login") or "").strip()
        topic = (body.get("topic") or "").strip()
        description = (body.get("description") or "").strip()
        audience = (body.get("audience") or "").strip()
        custom_design = bool(body.get("customDesign"))
        try:
            design_variant = int(body.get("designVariant") or 0)
        except (TypeError, ValueError):
            design_variant = 0

        try:
            slides_count = int(body.get("slidesCount") or 8)
        except (TypeError, ValueError):
            slides_count = 8
        slides_count = max(3, min(slides_count, 16))

        if not topic:
            return _resp(400, {"error": "Укажите тему урока"})

        # Проверяем лимит AI-запросов для trial-пользователей
        if login:
            try:
                limit_req = urllib.request.Request(
                    f"{AUTH_URL}?action=check-ai-limit",
                    data=json.dumps({"login": login}).encode("utf-8"),
                    method="POST",
                    headers={"Content-Type": "application/json"},
                )
                with urllib.request.urlopen(limit_req, timeout=10) as r:
                    limit_data = json.loads(r.read().decode())
                if not limit_data.get("allowed"):
                    return _resp(429, {"error": limit_data.get("error", "Достигнут лимит ИИ-запросов")})
            except urllib.error.HTTPError as e:
                err_body = json.loads(e.read().decode() or "{}")
                if e.code == 429:
                    return _resp(429, {"error": err_body.get("error", "Достигнут лимит ИИ-запросов на сегодня")})
            except Exception:
                pass

        # Предусматриваем расход: блокируем ИИ при отсутствии подписки/баланса.
        allowed, pc_status, pc_err = precheck_ai(login, est_tokens=4000)
        if not allowed:
            return _resp(pc_status, {"error": pc_err})

        try:
            outline, tokens_used = generate_outline(topic, description, slides_count, audience)
        except Exception as e:
            msg = str(e)
            if "timed out" in msg.lower() or "timeout" in msg.lower():
                return _resp(504, {"error": "ИИ-сервис не успел ответить — попробуйте ещё раз."})
            if "429" in msg or "rate" in msg.lower():
                return _resp(429, {"error": "Слишком много запросов к ИИ-сервису — подождите 30 секунд."})
            return _resp(500, {"error": f"Ошибка генерации структуры: {msg}"})

        _, _, spent_rub, balance_rub = spend_ai_tokens(login, max(tokens_used, 1))

        # Каждая презентация — уникальный современный дизайн (генерируется на лету).
        # Если variant не задан явно — берём случайный, чтобы одна тема давала разный стиль.
        variant = design_variant or random.randint(1, 10_000)
        theme = make_custom_theme(topic, variant)
        return _resp(200, {
            "outline": outline,
            "theme_name": theme["name"],
            "theme_payload": theme_to_payload(theme),
            "topic": topic,
            "spent_rub": spent_rub,
            "balance_rub": balance_rub,
        })

    # ── ШАГ 2: build — фото + PPTX (~15 сек) ───────────────────────────────
    if action == "build":
        topic = (body.get("topic") or "").strip()
        teacher_name = (body.get("teacherName") or "").strip()
        teacher_school = (body.get("teacherSchool") or "").strip()
        outline = body.get("outline")
        theme_name = (body.get("theme_name") or "").strip()
        theme_payload = body.get("theme_payload")
        regen_design = bool(body.get("regenDesign"))
        try:
            design_variant = int(body.get("designVariant") or 0)
        except (TypeError, ValueError):
            design_variant = 0

        if not topic or not outline:
            return _resp(400, {"error": "Укажите topic и outline"})

        # Восстанавливаем тему: при «обновить дизайн» генерируем новый индивидуальный
        # на лету, иначе из payload, иначе по имени.
        if regen_design:
            theme = make_custom_theme(topic, design_variant or random.randint(1, 10_000))
        elif theme_payload:
            theme = theme_from_payload(theme_payload)
        else:
            theme = make_custom_theme(topic, random.randint(1, 10_000))

        # Скачиваем до 3 фото на слайд параллельно
        images = {}
        try:
            img_list = fetch_images_parallel(outline.get("slides", []), topic)
            for idx, imgs in enumerate(img_list, start=1):
                if imgs:
                    images[idx] = imgs  # list[bytes]
        except Exception:
            pass

        try:
            pptx_bytes = build_pptx(
                topic=topic,
                subtitle=outline.get("subtitle", ""),
                contents=outline.get("contents", []),
                slides_data=outline.get("slides", []),
                conclusion=outline.get("conclusion", []),
                teacher_name=teacher_name,
                teacher_school=teacher_school,
                theme=theme,
                images=images,
            )
        except Exception as e:
            return _resp(500, {"error": f"Ошибка сборки PPTX: {e}"})

        filename = f"{safe_filename(topic)}.pptx"
        pptx_url = upload_pptx_to_s3(pptx_bytes, filename)
        resp_body = {
            "pptx_url": pptx_url,
            "filename": filename,
            "size": len(pptx_bytes),
            "outline": {
                "subtitle": outline.get("subtitle", ""),
                "slides": [{"title": s["title"], "bullets": s["bullets"]} for s in outline.get("slides", [])],
                "conclusion": outline.get("conclusion", []),
            },
        }
        # base64 кладём как запасной вариант только для небольших файлов (иначе 502)
        if not pptx_url or len(pptx_bytes) < 3_000_000:
            resp_body["pptx_b64"] = base64.b64encode(pptx_bytes).decode()
        return _resp(200, resp_body)

    # ── Legacy: оба шага вместе (action="" или не указан) ───────────────────
    login = (body.get("login") or "").strip()
    topic = (body.get("topic") or "").strip()
    description = (body.get("description") or "").strip()
    audience = (body.get("audience") or "").strip()
    teacher_name = (body.get("teacherName") or "").strip()
    teacher_school = (body.get("teacherSchool") or "").strip()
    custom_design = bool(body.get("customDesign"))

    try:
        slides_count = int(body.get("slidesCount") or 8)
    except (TypeError, ValueError):
        slides_count = 8
    slides_count = max(3, min(slides_count, 16))

    if not topic:
        return _resp(400, {"error": "Укажите тему урока"})

    if login:
        try:
            limit_req = urllib.request.Request(
                f"{AUTH_URL}?action=check-ai-limit",
                data=json.dumps({"login": login}).encode("utf-8"),
                method="POST",
                headers={"Content-Type": "application/json"},
            )
            with urllib.request.urlopen(limit_req, timeout=10) as r:
                limit_data = json.loads(r.read().decode())
            if not limit_data.get("allowed"):
                return _resp(429, {"error": limit_data.get("error", "Достигнут лимит ИИ-запросов")})
        except urllib.error.HTTPError as e:
            err_body = json.loads(e.read().decode() or "{}")
            if e.code == 429:
                return _resp(429, {"error": err_body.get("error", "Достигнут лимит ИИ-запросов на сегодня")})
        except Exception:
            pass

    theme = make_custom_theme(topic, random.randint(1, 10_000))

    # Предусматриваем расход: блокируем ИИ при отсутствии подписки/баланса.
    allowed, pc_status, pc_err = precheck_ai(login, est_tokens=4000)
    if not allowed:
        return _resp(pc_status, {"error": pc_err})

    try:
        outline, tokens_used = generate_outline(topic, description, slides_count, audience)
    except Exception as e:
        msg = str(e)
        if "timed out" in msg.lower() or "timeout" in msg.lower():
            return _resp(504, {"error": "ИИ-сервис не успел ответить — попробуйте ещё раз."})
        if "429" in msg or "rate" in msg.lower():
            return _resp(429, {"error": "Слишком много запросов к ИИ-сервису — подождите 30 секунд."})
        return _resp(500, {"error": f"Ошибка генерации структуры: {msg}"})

    _, _, spent_rub, balance_rub = spend_ai_tokens(login, max(tokens_used, 1))

    images = {}
    try:
        img_list = fetch_images_parallel(outline["slides"], topic)
        for idx, imgs in enumerate(img_list, start=1):
            if imgs:
                images[idx] = imgs  # list[bytes]
    except Exception:
        pass

    try:
        pptx_bytes = build_pptx(
            topic=topic,
            subtitle=outline["subtitle"],
            contents=outline.get("contents", []),
            slides_data=outline["slides"],
            conclusion=outline["conclusion"],
            teacher_name=teacher_name,
            teacher_school=teacher_school,
            theme=theme,
            images=images,
        )
    except Exception as e:
        return _resp(500, {"error": f"Ошибка сборки PPTX: {e}"})

    filename = f"{safe_filename(topic)}.pptx"
    pptx_url = upload_pptx_to_s3(pptx_bytes, filename)
    resp_body = {
        "pptx_url": pptx_url,
        "filename": filename,
        "size": len(pptx_bytes),
        "outline": {
            "subtitle": outline["subtitle"],
            "slides": [{"title": s["title"], "bullets": s["bullets"]} for s in outline["slides"]],
            "conclusion": outline["conclusion"],
        },
    }
    if not pptx_url or len(pptx_bytes) < 3_000_000:
        resp_body["pptx_b64"] = base64.b64encode(pptx_bytes).decode()
    return _resp(200, resp_body)